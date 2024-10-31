from pptx import Presentation
import json

def extract_slide_texts(ppt_path):
    presentation = Presentation(ppt_path)
    content = []
    
    for slide_number, slide in enumerate(presentation.slides, 1):
        slide_texts = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text_info = {
                        "text": paragraph.text.strip(),
                        "hyperlinks": []
                    }
                    
                    # Extract hyperlinks from runs
                    for run in paragraph.runs:
                        if run.hyperlink and run.hyperlink.address:
                            text_info["hyperlinks"].append({
                                "text": run.text,
                                "url": run.hyperlink.address
                            })
                    
                    if text_info["text"]:  # Only add non-empty texts
                        slide_texts.append(text_info)
        
        if slide_texts:  # Only add slides with content
            content.append({
                "slide": slide_number,
                "content": slide_texts
            })
    
    return content

if __name__ == "__main__":
    try:
        ppt_file = "sample.pptx"
        content = extract_slide_texts(ppt_file)
        
        # Save to JSON file
        with open("reversed_content.json", "w", encoding="utf-8") as f:
            json.dump(content, f, indent=2, ensure_ascii=False)
            
        print("Successfully extracted content to reversed_content.json")
        
    except FileNotFoundError:
        print(f"PowerPoint file '{ppt_file}' not found")
    except Exception as e:
        print(f"An error occurred: {str(e)}")